#!/usr/bin/env python3
"""Generate PowerPoint presentation of pLIN tool architecture."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
DARK_BLUE = RGBColor(0x0D, 0x47, 0xA1)
MED_BLUE = RGBColor(0x1E, 0x88, 0xE5)
LIGHT_BLUE = RGBColor(0xBB, 0xDE, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x75, 0x75, 0x75)
LIGHT_GRAY = RGBColor(0xEC, 0xEF, 0xF1)
RED = RGBColor(0xE5, 0x39, 0x35)
GREEN = RGBColor(0x43, 0xA0, 0x47)
ORANGE = RGBColor(0xFB, 0x8C, 0x00)
PURPLE = RGBColor(0x8E, 0x24, 0xAA)
TEAL = RGBColor(0x00, 0x96, 0x88)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=WHITE):
    """Set slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_box(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(1)):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=14, bold=False,
             color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=12,
                  color=DARK_GRAY, bold_first=False, spacing=1.0, font_name="Calibri"):
    """Add multi-line text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(spacing * 2)
        if bold_first and i == 0:
            p.font.bold = True
    return txBox


def add_arrow(slide, left, top, width, height, color=MED_BLUE):
    """Add a down arrow."""
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_right_arrow(slide, left, top, width, height, color=MED_BLUE):
    """Add a right arrow."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_header_bar(slide, text, subtitle=""):
    """Add a dark blue header bar at top."""
    add_box(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
    add_text(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.6),
             text, font_size=32, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(0.4),
                 subtitle, font_size=16, color=LIGHT_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)

add_text(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1),
         "pLIN — Plasmid Life Identification Number",
         font_size=40, bold=True, color=WHITE)

add_text(slide, Inches(0.8), Inches(2.6), Inches(11.5), Inches(0.8),
         "A Hierarchical, Reference-Free Classification System\nfor Bacterial Plasmid Genomes with AMR Surveillance",
         font_size=22, color=LIGHT_BLUE)

# Feature boxes
features = [
    ("6,346", "Training\nPlasmids"),
    ("3", "Inc Groups\n(FII, N, X1)"),
    ("2,232", "Unique pLIN\nCodes"),
    ("96.1%", "Inc Detection\nAccuracy"),
    ("27,465", "AMR Gene\nDetections"),
]
for i, (num, label) in enumerate(features):
    x = Inches(0.8 + i * 2.5)
    y = Inches(4.0)
    box = add_box(slide, x, y, Inches(2.2), Inches(1.8), RGBColor(0x0A, 0x2F, 0x6E),
                  border_color=MED_BLUE, border_width=Pt(2))
    add_text(slide, x, y + Inches(0.2), Inches(2.2), Inches(0.7),
             num, font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(0.9), Inches(2.2), Inches(0.7),
             label, font_size=14, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
         "Tool Architecture Overview", font_size=14, color=MED_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: System Overview / Architecture Diagram
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "System Architecture", "End-to-end pipeline from FASTA input to classified plasmids + AMR profile")

# Input
add_box(slide, Inches(0.4), Inches(1.6), Inches(2.4), Inches(1.4), LIGHT_BLUE, border_color=MED_BLUE, border_width=Pt(2))
add_text(slide, Inches(0.5), Inches(1.7), Inches(2.2), Inches(0.4),
         "INPUT", font_size=14, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
add_multiline(slide, Inches(0.5), Inches(2.1), Inches(2.2), Inches(0.8), [
    "Plasmid FASTA files",
    "(.fasta, .fa, .fna)",
], font_size=12, color=DARK_GRAY)

# Arrow 1
add_right_arrow(slide, Inches(2.95), Inches(2.0), Inches(0.6), Inches(0.4))

# Inc Detection
add_box(slide, Inches(3.7), Inches(1.6), Inches(2.6), Inches(1.4), RGBColor(0xFF, 0xF3, 0xE0), border_color=ORANGE, border_width=Pt(2))
add_text(slide, Inches(3.8), Inches(1.7), Inches(2.4), Inches(0.4),
         "INC GROUP DETECTION", font_size=12, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
add_multiline(slide, Inches(3.8), Inches(2.1), Inches(2.4), Inches(0.8), [
    "KNN Classifier (k=5)",
    "Cosine distance on 4-mers",
    "Trained on 6,346 plasmids",
], font_size=11, color=DARK_GRAY)

# Arrow 2
add_right_arrow(slide, Inches(6.45), Inches(2.0), Inches(0.6), Inches(0.4))

# K-mer + Clustering
add_box(slide, Inches(7.2), Inches(1.6), Inches(2.8), Inches(1.4), RGBColor(0xE8, 0xF5, 0xE9), border_color=GREEN, border_width=Pt(2))
add_text(slide, Inches(7.3), Inches(1.7), Inches(2.6), Inches(0.4),
         "pLIN ASSIGNMENT", font_size=12, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)
add_multiline(slide, Inches(7.3), Inches(2.1), Inches(2.6), Inches(0.8), [
    "4-mer frequency vectors (256D)",
    "Cosine distance matrix",
    "Single-linkage clustering",
    "6 threshold cuts \u2192 A.B.C.D.E.F",
], font_size=11, color=DARK_GRAY)

# Arrow 3
add_right_arrow(slide, Inches(10.15), Inches(2.0), Inches(0.6), Inches(0.4))

# Output
add_box(slide, Inches(10.9), Inches(1.6), Inches(2.1), Inches(1.4), RGBColor(0xE3, 0xF2, 0xFD), border_color=DARK_BLUE, border_width=Pt(2))
add_text(slide, Inches(11.0), Inches(1.7), Inches(1.9), Inches(0.4),
         "pLIN CODES", font_size=12, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
add_multiline(slide, Inches(11.0), Inches(2.1), Inches(1.9), Inches(0.8), [
    "Hierarchical codes",
    "e.g. 1.1.1.1.2.5",
    "Permanent & stable",
], font_size=11, color=DARK_GRAY)

# AMR Branch (parallel path)
add_box(slide, Inches(0.4), Inches(3.5), Inches(2.4), Inches(1.4), RGBColor(0xFC, 0xE4, 0xEC), border_color=RED, border_width=Pt(2))
add_text(slide, Inches(0.5), Inches(3.6), Inches(2.2), Inches(0.4),
         "AMRFinderPlus", font_size=14, bold=True, color=RED, alignment=PP_ALIGN.CENTER)
add_multiline(slide, Inches(0.5), Inches(4.0), Inches(2.2), Inches(0.8), [
    "NCBI AMR detection",
    "AMR + Stress + Virulence",
    "Per-contig gene calls",
], font_size=11, color=DARK_GRAY)

# Arrow from AMR
add_right_arrow(slide, Inches(2.95), Inches(3.9), Inches(0.6), Inches(0.4))

# Integration
add_box(slide, Inches(3.7), Inches(3.5), Inches(2.6), Inches(1.4), RGBColor(0xF3, 0xE5, 0xF5), border_color=PURPLE, border_width=Pt(2))
add_text(slide, Inches(3.8), Inches(3.6), Inches(2.4), Inches(0.4),
         "INTEGRATION", font_size=12, bold=True, color=PURPLE, alignment=PP_ALIGN.CENTER)
add_multiline(slide, Inches(3.8), Inches(4.0), Inches(2.4), Inches(0.8), [
    "pLIN codes + AMR genes",
    "Per-plasmid AMR profile",
    "Lineage AMR summary",
], font_size=11, color=DARK_GRAY)

# Arrow to Visualization
add_right_arrow(slide, Inches(6.45), Inches(3.9), Inches(0.6), Inches(0.4))

# Visualization
add_box(slide, Inches(7.2), Inches(3.5), Inches(2.8), Inches(1.4), RGBColor(0xE0, 0xF7, 0xFA), border_color=TEAL, border_width=Pt(2))
add_text(slide, Inches(7.3), Inches(3.6), Inches(2.6), Inches(0.4),
         "VISUALIZATION", font_size=12, bold=True, color=TEAL, alignment=PP_ALIGN.CENTER)
add_multiline(slide, Inches(7.3), Inches(4.0), Inches(2.6), Inches(0.8), [
    "Cladograms (4 types)",
    "AMR heatmaps (Plotly)",
    "Drug class pie charts",
    "Critical gene alerts",
], font_size=11, color=DARK_GRAY)

# Arrow to Export
add_right_arrow(slide, Inches(10.15), Inches(3.9), Inches(0.6), Inches(0.4))

# Export
add_box(slide, Inches(10.9), Inches(3.5), Inches(2.1), Inches(1.4), LIGHT_GRAY, border_color=MED_GRAY, border_width=Pt(2))
add_text(slide, Inches(11.0), Inches(3.6), Inches(1.9), Inches(0.4),
         "EXPORT", font_size=12, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_multiline(slide, Inches(11.0), Inches(4.0), Inches(1.9), Inches(0.8), [
    "TSV tables",
    "PNG / PDF figures",
    "ZIP bundle",
], font_size=11, color=DARK_GRAY)

# GUI bar at bottom
add_box(slide, Inches(0.4), Inches(5.5), Inches(12.6), Inches(1.5), RGBColor(0xE8, 0xEA, 0xF6), border_color=RGBColor(0x53, 0x4B, 0xAE), border_width=Pt(2))
add_text(slide, Inches(0.6), Inches(5.6), Inches(3), Inches(0.4),
         "Streamlit Web GUI (plin_app.py)", font_size=14, bold=True, color=RGBColor(0x53, 0x4B, 0xAE))
tabs = ["Upload & Config", "Overview Tab", "Results Tab", "Cladogram Tab", "AMR Analysis Tab", "Export Tab"]
for i, tab in enumerate(tabs):
    x = Inches(0.6 + i * 2.05)
    add_box(slide, x, Inches(6.1), Inches(1.9), Inches(0.6), WHITE, border_color=RGBColor(0x53, 0x4B, 0xAE))
    add_text(slide, x, Inches(6.15), Inches(1.9), Inches(0.5),
             tab, font_size=11, bold=True, color=RGBColor(0x53, 0x4B, 0xAE), alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: pLIN Classification System
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "pLIN Classification System", "Six-level hierarchical coding from tetranucleotide composition")

# Threshold table
levels = [
    ("A", "Family", "0.150", "~85%", "#E53935"),
    ("B", "Subfamily", "0.100", "~90%", "#FB8C00"),
    ("C", "Cluster", "0.050", "~95%", "#FDD835"),
    ("D", "Subcluster", "0.020", "~98%", "#43A047"),
    ("E", "Clone", "0.010", "~99%", "#1E88E5"),
    ("F", "Strain", "0.001", "~99.9%", "#8E24AA"),
]

# Table header
add_box(slide, Inches(0.5), Inches(1.6), Inches(5.5), Inches(0.5), DARK_BLUE)
headers = ["Position", "Level", "Threshold (d\u2264)", "ANI Equiv."]
for i, h in enumerate(headers):
    add_text(slide, Inches(0.6 + i * 1.35), Inches(1.63), Inches(1.3), Inches(0.4),
             h, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

for j, (pos, level, thresh, ani, hex_col) in enumerate(levels):
    y = Inches(2.1 + j * 0.45)
    bg = LIGHT_GRAY if j % 2 == 0 else WHITE
    add_box(slide, Inches(0.5), y, Inches(5.5), Inches(0.45), bg)
    r, g, b = int(hex_col[1:3], 16), int(hex_col[3:5], 16), int(hex_col[5:7], 16)
    # Color dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.1), Inches(0.25), Inches(0.25))
    dot.fill.solid()
    dot.fill.fore_color.rgb = RGBColor(r, g, b)
    dot.line.fill.background()
    add_text(slide, Inches(1.0), y + Inches(0.05), Inches(0.9), Inches(0.35),
             pos, font_size=13, bold=True, color=RGBColor(r, g, b), alignment=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.95), y + Inches(0.05), Inches(1.3), Inches(0.35),
             level, font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    add_text(slide, Inches(3.3), y + Inches(0.05), Inches(1.3), Inches(0.35),
             f"d \u2264 {thresh}", font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER, font_name="Consolas")
    add_text(slide, Inches(4.65), y + Inches(0.05), Inches(1.3), Inches(0.35),
             ani, font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Example pLIN code
add_box(slide, Inches(0.5), Inches(5.0), Inches(5.5), Inches(1.0), RGBColor(0xE3, 0xF2, 0xFD), border_color=MED_BLUE, border_width=Pt(2))
add_text(slide, Inches(0.7), Inches(5.1), Inches(5), Inches(0.4),
         "Example pLIN Code:", font_size=13, bold=True, color=DARK_BLUE)
add_text(slide, Inches(0.7), Inches(5.5), Inches(5), Inches(0.4),
         "1 . 1 . 1 . 1 . 2 . 5", font_size=24, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER, font_name="Consolas")

# Right side: Algorithm steps
add_text(slide, Inches(6.8), Inches(1.6), Inches(6), Inches(0.5),
         "Classification Algorithm", font_size=20, bold=True, color=DARK_BLUE)

steps = [
    ("1", "Parse FASTA", "Read plasmid sequences using BioPython", GREEN),
    ("2", "Compute 4-mer Vectors", "256-dimensional normalised frequency vectors", ORANGE),
    ("3", "Cosine Distance Matrix", "Pairwise distance between all sequences", MED_BLUE),
    ("4", "Single-Linkage Clustering", "Hierarchical agglomerative clustering (scipy)", PURPLE),
    ("5", "Threshold Cutting", "Cut dendrogram at 6 thresholds \u2192 6 cluster IDs", RED),
    ("6", "Assign pLIN Codes", "Concatenate cluster IDs: A.B.C.D.E.F", DARK_BLUE),
]

for i, (num, title, desc, color) in enumerate(steps):
    y = Inches(2.2 + i * 0.72)
    # Number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.8), y, Inches(0.4), Inches(0.4))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    tf = circ.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text(slide, Inches(7.4), y - Inches(0.02), Inches(5.5), Inches(0.3),
             title, font_size=14, bold=True, color=color)
    add_text(slide, Inches(7.4), y + Inches(0.28), Inches(5.5), Inches(0.3),
             desc, font_size=11, color=MED_GRAY)

    if i < len(steps) - 1:
        # Connecting line
        line = slide.shapes.add_connector(1, Inches(7.0), y + Inches(0.4), Inches(7.0), y + Inches(0.72))
        line.line.color.rgb = LIGHT_BLUE
        line.line.width = Pt(2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: Inc Group Auto-Detection
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Inc Group Auto-Detection", "KNN classifier trained on 6,346 RefSeq plasmids")

# Training data
add_text(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(0.5),
         "Training Dataset", font_size=20, bold=True, color=DARK_BLUE)

groups = [
    ("IncFII", "4,581", "72.2%", RED),
    ("IncN", "1,064", "16.8%", GREEN),
    ("IncX1", "701", "11.0%", MED_BLUE),
]

for i, (name, count, pct, color) in enumerate(groups):
    y = Inches(2.2 + i * 0.9)
    bar_width = float(pct.replace("%", "")) / 100 * 5.5
    add_box(slide, Inches(0.5), y, Inches(bar_width), Inches(0.6), color)
    add_text(slide, Inches(0.7), y + Inches(0.1), Inches(3), Inches(0.4),
             f"{name}  —  {count} plasmids ({pct})", font_size=14, bold=True, color=WHITE)

# Classifier details
add_text(slide, Inches(0.5), Inches(5.0), Inches(6), Inches(0.5),
         "Classifier Details", font_size=18, bold=True, color=DARK_BLUE)
details = [
    "Algorithm: K-Nearest Neighbors (k=5, distance-weighted)",
    "Distance metric: Cosine distance on 4-mer frequency vectors",
    "5-fold cross-validation accuracy: 96.1% \u00b1 0.6%",
    "Hold-out test accuracy: 100% (200 samples per group)",
    "Model file: data/inc_classifier.npz (4.3 MB)",
]
add_multiline(slide, Inches(0.5), Inches(5.5), Inches(6), Inches(1.5), details, font_size=13, color=DARK_GRAY)

# Right side: How it works
add_text(slide, Inches(7.0), Inches(1.6), Inches(6), Inches(0.5),
         "How Auto-Detection Works", font_size=20, bold=True, color=DARK_BLUE)

steps_inc = [
    "Upload new plasmid FASTA",
    "Compute 4-mer frequency vector (256 features)",
    "Find 5 nearest neighbors in training set (cosine distance)",
    "Distance-weighted vote \u2192 Inc group prediction",
    "Return predicted group + confidence probability",
]
for i, step in enumerate(steps_inc):
    y = Inches(2.3 + i * 0.65)
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.0), y, Inches(0.35), Inches(0.35))
    circ.fill.solid()
    circ.fill.fore_color.rgb = ORANGE
    circ.line.fill.background()
    tf = circ.text_frame
    tf.paragraphs[0].text = str(i + 1)
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text(slide, Inches(7.5), y, Inches(5.5), Inches(0.35),
             step, font_size=13, color=DARK_GRAY)

# Output example
add_box(slide, Inches(7.0), Inches(5.7), Inches(5.8), Inches(1.3), RGBColor(0xFf, 0xF8, 0xE1), border_color=ORANGE, border_width=Pt(2))
add_text(slide, Inches(7.2), Inches(5.8), Inches(5.4), Inches(0.3),
         "Example Output:", font_size=12, bold=True, color=ORANGE)
add_multiline(slide, Inches(7.2), Inches(6.1), Inches(5.4), Inches(0.8), [
    "Plasmid: SP12_P2  \u2192  IncX1 (confidence: 0.79)",
    "  IncFII: 0.00  |  IncN: 0.21  |  IncX1: 0.79",
], font_size=12, color=DARK_GRAY, font_name="Consolas")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: AMRFinderPlus Integration
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "AMRFinderPlus Integration", "NCBI antimicrobial resistance gene detection pipeline")

# Left: AMR pipeline
add_text(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(0.5),
         "AMR Detection Pipeline", font_size=20, bold=True, color=RED)

amr_steps = [
    ("Auto-Detect", "Search PATH, conda envs for amrfinder binary", RGBColor(0x78, 0x90, 0x9C)),
    ("Run per File", "amrfinder -n {fasta} -d {db} --plus -o {output}", RED),
    ("Parse Results", "Read TSV output: gene, class, type, coordinates", ORANGE),
    ("Integrate", "Merge AMR hits with pLIN assignments per plasmid", PURPLE),
    ("Visualize", "Heatmaps, bar charts, pie charts, alerts", TEAL),
]

for i, (title, desc, color) in enumerate(amr_steps):
    y = Inches(2.3 + i * 0.85)
    add_box(slide, Inches(0.5), y, Inches(5.8), Inches(0.7), WHITE, border_color=color, border_width=Pt(2))
    add_text(slide, Inches(0.7), y + Inches(0.05), Inches(1.8), Inches(0.3),
             title, font_size=13, bold=True, color=color)
    add_text(slide, Inches(0.7), y + Inches(0.35), Inches(5.4), Inches(0.3),
             desc, font_size=11, color=MED_GRAY)

# Right: Detection types
add_text(slide, Inches(7.0), Inches(1.6), Inches(6), Inches(0.5),
         "Detection Categories", font_size=20, bold=True, color=RED)

categories = [
    ("AMR Genes", "Antibiotic resistance genes\nblaNDM, blaKPC, blaSHV, mcr, etc.", RED, "27,465"),
    ("Stress Genes", "Biocide/metal resistance\nbleomycin, mercury, arsenic, etc.", ORANGE, "5,834"),
    ("Virulence", "Virulence-associated factors\nType III secretion, toxins, etc.", PURPLE, "1,200+"),
]

for i, (title, desc, color, count) in enumerate(categories):
    y = Inches(2.3 + i * 1.4)
    add_box(slide, Inches(7.0), y, Inches(5.8), Inches(1.2), WHITE, border_color=color, border_width=Pt(2))
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.3), y + Inches(0.35), Inches(0.4), Inches(0.4))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_text(slide, Inches(7.9), y + Inches(0.1), Inches(3), Inches(0.35),
             title, font_size=15, bold=True, color=color)
    add_text(slide, Inches(7.9), y + Inches(0.5), Inches(4.5), Inches(0.6),
             desc, font_size=11, color=MED_GRAY)
    add_text(slide, Inches(11.0), y + Inches(0.2), Inches(1.5), Inches(0.8),
             count, font_size=22, bold=True, color=color, alignment=PP_ALIGN.RIGHT)

# Critical gene alert box
add_box(slide, Inches(7.0), Inches(6.0), Inches(5.8), Inches(1.0), RGBColor(0xFF, 0xEB, 0xEE), border_color=RED, border_width=Pt(2))
add_text(slide, Inches(7.2), Inches(6.1), Inches(5.4), Inches(0.3),
         "\u26a0\ufe0f  Critical Gene Alert System", font_size=13, bold=True, color=RED)
add_text(slide, Inches(7.2), Inches(6.45), Inches(5.4), Inches(0.4),
         "Carbapenemases (KPC, NDM, OXA-48, VIM, IMP)  |  ESBLs (CTX-M, SHV-12, TEM)  |  Colistin (mcr)",
         font_size=11, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: Visualization & GUI
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Streamlit Web GUI", "Interactive web application for plasmid analysis (plin_app.py)")

# GUI Layout diagram
add_text(slide, Inches(0.5), Inches(1.6), Inches(12), Inches(0.5),
         "Application Layout", font_size=20, bold=True, color=DARK_BLUE)

# Upload area
add_box(slide, Inches(0.5), Inches(2.2), Inches(8.5), Inches(1.2), LIGHT_BLUE, border_color=MED_BLUE, border_width=Pt(2))
add_text(slide, Inches(0.7), Inches(2.3), Inches(4), Inches(0.4),
         "\u2191 Upload FASTA Files", font_size=14, bold=True, color=DARK_BLUE)
add_text(slide, Inches(5.5), Inches(2.3), Inches(3), Inches(0.4),
         "Inc Group: [Auto-detect \u25bc]", font_size=12, color=DARK_GRAY)
add_text(slide, Inches(5.5), Inches(2.7), Inches(3), Inches(0.4),
         "\u2611 Run AMRFinderPlus", font_size=12, color=DARK_GRAY)
add_box(slide, Inches(0.7), Inches(2.9), Inches(2), Inches(0.4), MED_BLUE)
add_text(slide, Inches(0.7), Inches(2.92), Inches(2), Inches(0.35),
         "\u25b6 Run Analysis", font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Tabs
tabs_info = [
    ("Overview", "pLIN description\nThreshold table\nMetrics dashboard\nAdaptive threshold info", MED_BLUE),
    ("Results", "Interactive data table\nSearch/filter\nMobility column\npLIN distribution", GREEN),
    ("Cladogram", "4 visualization types:\n\u2022 Rectangular\n\u2022 Circular\n\u2022 Heatmap\n\u2022 AMR Annotated", ORANGE),
    ("AMR", "Gene prevalence\nDrug class pie\nCritical gene alerts\nFull detection table", RED),
    ("Epidemiology", "Mobility prediction\nOutbreak detection\nRisk stratification\nDissemination risk", TEAL),
    ("Export", "TSV tables\nPNG/PDF figures\nZIP bundle\nMobility + outbreak", PURPLE),
]

for i, (name, desc, color) in enumerate(tabs_info):
    x = Inches(0.5 + i * 2.1)
    y = Inches(3.8)
    add_box(slide, x, y, Inches(1.95), Inches(0.45), color)
    add_text(slide, x, y + Inches(0.05), Inches(1.95), Inches(0.35),
             name, font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_box(slide, x, y + Inches(0.45), Inches(1.95), Inches(2.5), WHITE, border_color=color, border_width=Pt(1))
    add_multiline(slide, x + Inches(0.05), y + Inches(0.55), Inches(1.85), Inches(2.3),
                  desc.split("\n"), font_size=9, color=DARK_GRAY)

# Sidebar
add_box(slide, Inches(9.3), Inches(2.2), Inches(3.7), Inches(4.8), LIGHT_GRAY, border_color=MED_GRAY, border_width=Pt(1))
add_text(slide, Inches(9.5), Inches(2.3), Inches(3.3), Inches(0.4),
         "Sidebar", font_size=16, bold=True, color=DARK_GRAY)
sidebar_items = [
    "Analysis status & metrics",
    "Plasmid count",
    "Unique pLIN codes",
    "Detected Inc groups",
    "",
    "\U0001f504 Clear & Reset",
    "\U0001f4c2 New Analysis",
]
add_multiline(slide, Inches(9.5), Inches(2.8), Inches(3.3), Inches(4), sidebar_items, font_size=11, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: File Structure & Project Organization
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Project Structure", "File organization and data flow")

# Left column: Scripts
add_text(slide, Inches(0.5), Inches(1.6), Inches(4), Inches(0.5),
         "Python Scripts", font_size=18, bold=True, color=DARK_BLUE)

scripts = [
    ("plin_app.py", "Streamlit GUI (main app)", MED_BLUE),
    ("assign_pLIN.py", "Batch pLIN assignment", GREEN),
    ("integrate_pLIN_AMR.py", "pLIN + AMR integration", PURPLE),
    ("generate_figures.py", "Publication figures", ORANGE),
    ("build_inc_centroids.py", "Train Inc classifier", RED),
    ("test_pLIN.py", "Test pLIN on 22 plasmids", MED_GRAY),
    ("test_cladogram.py", "Test cladogram generation", MED_GRAY),
    ("test_integrate_and_cladogram.py", "Test AMR + cladogram", MED_GRAY),
]

for i, (fname, desc, color) in enumerate(scripts):
    y = Inches(2.1 + i * 0.46)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y + Inches(0.08), Inches(0.2), Inches(0.2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_text(slide, Inches(0.9), y, Inches(2), Inches(0.35),
             fname, font_size=11, bold=True, color=DARK_GRAY, font_name="Consolas")
    add_text(slide, Inches(3.0), y, Inches(2), Inches(0.35),
             desc, font_size=10, color=MED_GRAY)

# Middle column: Data
add_text(slide, Inches(5.2), Inches(1.6), Inches(3.5), Inches(0.5),
         "Data & Training", font_size=18, bold=True, color=DARK_BLUE)

data_items = [
    "plasmid_sequences_for_training/",
    "  \u251c\u2500 IncFII/fastas/  (4,581 files)",
    "  \u251c\u2500 IncN/fastas/    (1,064 files)",
    "  \u2514\u2500 IncX1/fastas/   (701 files)",
    "",
    "data/",
    "  \u251c\u2500 inc_classifier.npz (4.3 MB)",
    "  \u2514\u2500 inc_centroids.npz",
    "",
    "test_plasmids/",
    "  \u2514\u2500 IncX/  (22 FASTA files)",
]

add_multiline(slide, Inches(5.2), Inches(2.1), Inches(3.8), Inches(4.5),
              data_items, font_size=11, color=DARK_GRAY, font_name="Consolas")

# Right column: Output
add_text(slide, Inches(9.5), Inches(1.6), Inches(3.5), Inches(0.5),
         "Output", font_size=18, bold=True, color=DARK_BLUE)

output_items = [
    "output/",
    "  \u251c\u2500 pLIN_assignments.tsv",
    "  \u251c\u2500 amrfinder/",
    "  \u2502   \u2514\u2500 amrfinder_all_plasmids.tsv",
    "  \u251c\u2500 integrated/",
    "  \u2502   \u251c\u2500 pLIN_AMR_integrated.tsv",
    "  \u2502   \u2514\u2500 pLIN_lineage_AMR_summary.tsv",
    "  \u251c\u2500 figures/ (7 figs, PNG+PDF)",
    "  \u2514\u2500 test/ (test results)",
]

add_multiline(slide, Inches(9.5), Inches(2.1), Inches(3.5), Inches(4.5),
              output_items, font_size=11, color=DARK_GRAY, font_name="Consolas")

# Bottom: Shell scripts
add_box(slide, Inches(0.5), Inches(6.0), Inches(12.4), Inches(1.0), LIGHT_GRAY, border_color=MED_GRAY)
add_text(slide, Inches(0.7), Inches(6.05), Inches(3), Inches(0.35),
         "Automation Scripts:", font_size=13, bold=True, color=DARK_GRAY)
shell_items = [
    "setup.sh / setup.bat  \u2014  Create venv & install deps",
    "run_all.sh / run_all.bat  \u2014  Full pipeline: pLIN \u2192 AMR \u2192 Integration \u2192 Figures",
    "run_amrfinder_all.sh  \u2014  Batch AMRFinderPlus on all training data",
]
add_multiline(slide, Inches(0.7), Inches(6.35), Inches(11.5), Inches(0.6),
              shell_items, font_size=11, color=DARK_GRAY, font_name="Consolas")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: Technology Stack
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Technology Stack", "Core dependencies and tools")

categories_tech = [
    ("Core Scientific", [
        ("NumPy \u2265 2.0", "Array operations, k-mer vectors"),
        ("Pandas \u2265 2.0", "DataFrames, TSV I/O"),
        ("SciPy \u2265 1.12", "pdist, linkage, fcluster, dendrogram"),
        ("BioPython \u2265 1.80", "FASTA parsing (SeqIO)"),
    ], MED_BLUE),
    ("Machine Learning", [
        ("scikit-learn \u2265 1.4", "KNN classifier for Inc detection"),
        ("XGBoost \u2265 2.0", "Validation classifier (F1=0.903)"),
        ("Optuna \u2265 3.5", "Hyperparameter optimization"),
    ], GREEN),
    ("Visualization", [
        ("Matplotlib \u2265 3.8", "Cladograms, static figures"),
        ("Seaborn \u2265 0.13", "Statistical plots, heatmaps"),
        ("Plotly \u2265 5.18", "Interactive charts (AMR tab)"),
    ], ORANGE),
    ("Web & External", [
        ("Streamlit \u2265 1.31", "Web GUI framework"),
        ("AMRFinderPlus", "NCBI AMR gene detection"),
        ("python-pptx", "PowerPoint generation"),
    ], PURPLE),
]

for i, (cat_name, items, color) in enumerate(categories_tech):
    x = Inches(0.5 + i * 3.2)
    add_box(slide, x, Inches(1.6), Inches(3.0), Inches(0.5), color)
    add_text(slide, x, Inches(1.63), Inches(3.0), Inches(0.4),
             cat_name, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    for j, (lib, desc) in enumerate(items):
        y = Inches(2.3 + j * 0.75)
        add_box(slide, x, y, Inches(3.0), Inches(0.65), WHITE, border_color=color, border_width=Pt(1))
        add_text(slide, x + Inches(0.1), y + Inches(0.03), Inches(2.8), Inches(0.3),
                 lib, font_size=12, bold=True, color=color, font_name="Consolas")
        add_text(slide, x + Inches(0.1), y + Inches(0.33), Inches(2.8), Inches(0.3),
                 desc, font_size=10, color=MED_GRAY)

# Key metrics at bottom
add_box(slide, Inches(0.5), Inches(5.6), Inches(12.4), Inches(1.4), RGBColor(0xE8, 0xF5, 0xE9), border_color=GREEN, border_width=Pt(2))
add_text(slide, Inches(0.7), Inches(5.7), Inches(12), Inches(0.4),
         "Key Performance Metrics", font_size=16, bold=True, color=GREEN)

metrics = [
    ("Simpson's D", "0.979", "Discriminatory power of pLIN system"),
    ("CV Accuracy", "96.1%", "Inc group auto-detection (5-fold)"),
    ("XGBoost F1", "0.903", "ML validation of k-mer features"),
    ("Concordance", "99.5%", "Composition vs known Inc groups"),
]

for i, (name, value, desc) in enumerate(metrics):
    x = Inches(0.7 + i * 3.1)
    add_text(slide, x, Inches(6.1), Inches(2.8), Inches(0.35),
             f"{name}: {value}", font_size=14, bold=True, color=DARK_BLUE)
    add_text(slide, x, Inches(6.45), Inches(2.8), Inches(0.3),
             desc, font_size=11, color=MED_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: Advanced Features (NEW)
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Advanced Analytical Features", "Epidemiological intelligence integrated into the pLIN pipeline")

# Feature 1: Adaptive Thresholds
add_box(slide, Inches(0.4), Inches(1.5), Inches(6.2), Inches(2.5), WHITE, border_color=MED_BLUE, border_width=Pt(2))
add_text(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(0.4),
         "Per-Inc-Group Adaptive Thresholds", font_size=16, bold=True, color=MED_BLUE)
add_multiline(slide, Inches(0.6), Inches(2.1), Inches(5.8), Inches(1.6), [
    "Calibrates pLIN distance thresholds per Inc group from training data",
    "Uses quantile-based calibration on within-group distance distributions",
    "Addresses the limitation that fixed thresholds may not suit all Inc groups",
    "Example: IncFII (diverse, 4,581 plasmids) needs wider thresholds than IncX1 (compact, 701)",
    "Automatically selects thresholds for the dominant Inc group in each analysis",
], font_size=11, color=DARK_GRAY)

# Feature 2: Linkage Method Selection
add_box(slide, Inches(6.9), Inches(1.5), Inches(6.0), Inches(2.5), WHITE, border_color=GREEN, border_width=Pt(2))
add_text(slide, Inches(7.1), Inches(1.6), Inches(5.6), Inches(0.4),
         "Multi-Linkage Clustering", font_size=16, bold=True, color=GREEN)
add_multiline(slide, Inches(7.1), Inches(2.1), Inches(5.6), Inches(1.6), [
    "User-selectable linkage method: Single, Complete, Average, Weighted",
    "Single (default): traditional chaining, sensitive to intermediates",
    "Complete: maximum inter-cluster distance, produces tighter clusters",
    "Average (UPGMA): balanced, widely used in phylogenetics",
    "Reduces chaining artifacts that can merge distinct lineages",
], font_size=11, color=DARK_GRAY)

# Feature 3: Mobility Prediction
add_box(slide, Inches(0.4), Inches(4.3), Inches(6.2), Inches(2.8), WHITE, border_color=ORANGE, border_width=Pt(2))
add_text(slide, Inches(0.6), Inches(4.4), Inches(5.8), Inches(0.4),
         "Plasmid Mobility Prediction", font_size=16, bold=True, color=ORANGE)
add_multiline(slide, Inches(0.6), Inches(4.9), Inches(5.8), Inches(2.0), [
    "Classifies plasmids as Conjugative, Mobilizable, or Non-mobilizable",
    "Scans AMRFinderPlus output for transfer/mobilization gene markers:",
    "  Conjugative: tra/trb genes (Type IV secretion system)",
    "  Mobilizable: mob genes (relaxase, oriT)",
    "Conjugative + AMR plasmids flagged as HIGH RISK for dissemination",
    "Enables risk-stratified surveillance of AMR-carrying plasmids",
], font_size=11, color=DARK_GRAY)

# Feature 4: Outbreak Detection
add_box(slide, Inches(6.9), Inches(4.3), Inches(6.0), Inches(2.8), WHITE, border_color=RED, border_width=Pt(2))
add_text(slide, Inches(7.1), Inches(4.4), Inches(5.6), Inches(0.4),
         "Outbreak / Clone Detection", font_size=16, bold=True, color=RED)
add_multiline(slide, Inches(7.1), Inches(4.9), Inches(5.6), Inches(2.0), [
    "Flags suspected outbreak clusters automatically",
    "Detection criteria: same pLIN strain code (F-level) + identical AMR profile",
    "Risk levels: HIGH (3+ shared AMR genes) / MODERATE (1-2 shared genes)",
    "Supports real-time surveillance for plasmid-mediated AMR spread",
    "Integrated into new Epidemiology tab with dissemination risk matrix",
    "Exportable as JSON for downstream analysis pipelines",
], font_size=11, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10: Novelty
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Novelty of pLIN System", "What distinguishes pLIN from existing plasmid classification approaches")

# Left column: Novelty points
novelty_items = [
    ("Reference-Free Classification",
     "Unlike MOB-suite, PlasmidFinder, or replicon typing that rely on curated reference databases, "
     "pLIN uses intrinsic sequence composition (4-mer frequencies). New or divergent plasmids that lack "
     "known replicons are still classifiable — no database gaps.",
     MED_BLUE),
    ("Permanent, Hierarchical Codes",
     "pLIN codes are stable and never change when new sequences are added. Existing plasmid typing schemes "
     "(e.g., pMLST) can reassign types as databases grow. The six-level hierarchy (Family \u2192 Strain) provides "
     "resolution at multiple scales in a single code — no existing system offers this.",
     GREEN),
    ("Composition-Based Inc Group Detection",
     "Most Inc group detection tools (PlasmidFinder, COPLA) require gene-level BLAST searches. pLIN's KNN "
     "classifier predicts Inc group from global k-mer composition alone (96% accuracy), enabling classification "
     "even when replicon genes are fragmented, truncated, or absent from assemblies.",
     ORANGE),
    ("Integrated AMR Surveillance Pipeline",
     "pLIN is the first system to natively integrate plasmid hierarchical classification with AMRFinderPlus "
     "gene detection in a single automated pipeline. This enables direct mapping of resistance gene cargo "
     "onto the plasmid phylogeny — linking plasmid backbone evolution with AMR gene acquisition.",
     RED),
    ("Interactive GUI with Epidemiological Intelligence",
     "Unlike command-line-only tools (MOB-suite, plasmidfinder CLI), pLIN ships a Streamlit web app with "
     "upload, classification, AMR screening, cladogram visualization, mobility prediction, outbreak detection, "
     "adaptive thresholds, and multi-linkage clustering — accessible to non-bioinformaticians.",
     PURPLE),
]

for i, (title, desc, color) in enumerate(novelty_items):
    y = Inches(1.5 + i * 1.15)
    # Colored left bar
    add_box(slide, Inches(0.4), y, Inches(0.12), Inches(0.95), color)
    add_text(slide, Inches(0.7), y + Inches(0.02), Inches(12.2), Inches(0.35),
             title, font_size=15, bold=True, color=color)
    add_text(slide, Inches(0.7), y + Inches(0.38), Inches(12.2), Inches(0.55),
             desc, font_size=11, color=DARK_GRAY)

# Comparison table at bottom
add_box(slide, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.35), DARK_BLUE)
comp_headers = ["Feature", "pLIN", "MOB-suite", "PlasmidFinder", "pMLST"]
col_widths = [2.8, 2.2, 2.5, 2.5, 2.5]
x_pos = 0.4
for j, (h, w) in enumerate(zip(comp_headers, col_widths)):
    add_text(slide, Inches(x_pos + 0.05), Inches(7.02), Inches(w), Inches(0.3),
             h, font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    x_pos += w


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11: Limitations & Future Directions
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide, "Limitations & Future Directions", "Current constraints and planned improvements")

# Limitations (left side)
add_text(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(0.5),
         "Current Limitations", font_size=20, bold=True, color=RED)

limitations = [
    ("Limited Inc Group Coverage",
     "Classifier trained on 3 Inc groups only (IncFII, IncN, IncX1). Plasmids from IncA/C, IncI, "
     "IncL/M, IncP, IncQ, ColE-type, and other groups will be misclassified into the nearest "
     "available category. Expanding training data is essential."),
    ("Composition-Only Features",
     "4-mer frequency captures global sequence composition but ignores gene content, synteny, "
     "and structural rearrangements. Two plasmids with similar base composition but different "
     "gene cargo may receive similar pLIN codes at coarse levels."),
    ("Linkage Method Sensitivity",
     "Single-linkage clustering (default) can produce chain-like clusters. While the tool now supports "
     "complete, average, and weighted linkage as alternatives, the pLIN thresholds were originally "
     "calibrated for single-linkage — using other methods may require re-calibration."),
    ("No Fragmented Assembly Handling",
     "pLIN expects complete or near-complete plasmid sequences. Short contigs from fragmented "
     "assemblies will have noisy k-mer profiles, reducing classification accuracy. No scaffolding "
     "or multi-contig plasmid reconstruction is performed."),
    ("Threshold Calibration Coverage",
     "Adaptive per-Inc-group thresholds are now available, calibrated from training data quantiles. "
     "However, calibration quality depends on training set size and diversity — smaller groups "
     "(IncX1: 701 plasmids) may have less reliable thresholds than larger ones (IncFII: 4,581)."),
    ("No Real-Time Database Updates",
     "The KNN classifier uses a static training set. New plasmid submissions to GenBank/RefSeq "
     "are not automatically incorporated. Periodic retraining is needed to maintain accuracy "
     "as novel plasmid lineages emerge."),
]

for i, (title, desc) in enumerate(limitations):
    y = Inches(2.15 + i * 0.82)
    # Red number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), y + Inches(0.05), Inches(0.3), Inches(0.3))
    circ.fill.solid()
    circ.fill.fore_color.rgb = RED
    circ.line.fill.background()
    tf = circ.text_frame
    tf.paragraphs[0].text = str(i + 1)
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text(slide, Inches(0.95), y, Inches(5.8), Inches(0.3),
             title, font_size=12, bold=True, color=RED)
    add_text(slide, Inches(0.95), y + Inches(0.3), Inches(5.8), Inches(0.5),
             desc, font_size=9, color=DARK_GRAY)

# Future Directions (right side)
add_text(slide, Inches(7.0), Inches(1.6), Inches(6), Inches(0.5),
         "Future Directions", font_size=20, bold=True, color=GREEN)

futures = [
    ("Expand Inc Group Training",
     "Add IncA/C, IncI, IncL/M, IncP, IncQ, ColE, and non-typeable plasmids. "
     "Target: 15+ Inc groups, 20,000+ training sequences.",
     GREEN),
    ("Hybrid Features",
     "Combine k-mer composition with gene presence/absence and synteny features "
     "for higher-resolution classification at fine levels (E, F).",
     MED_BLUE),
    ("Threshold Cross-Validation",
     "Systematically validate adaptive thresholds across Inc groups. "
     "Benchmark linkage methods against known plasmid phylogenies.",
     ORANGE),
    ("Metagenomic Support",
     "Handle multi-contig plasmid bins from metagenomic assemblies. "
     "Integrate with plasmid prediction tools (PlasFlow, MOB-recon).",
     PURPLE),
    ("Online Database & API",
     "Web database for pLIN code lookups and new sequence submission. "
     "REST API for programmatic access and LIMS integration.",
     TEAL),
    ("Multi-Hospital Outbreak Tracking",
     "Extend outbreak detection to cross-institutional datasets. "
     "Integrate with epidemiological metadata for spatial-temporal tracking.",
     RED),
]

for i, (title, desc, color) in enumerate(futures):
    y = Inches(2.15 + i * 0.82)
    # Arrow shape
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7.0), y + Inches(0.1), Inches(0.35), Inches(0.25))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()

    add_text(slide, Inches(7.5), y, Inches(5.3), Inches(0.3),
             title, font_size=12, bold=True, color=color)
    add_text(slide, Inches(7.5), y + Inches(0.3), Inches(5.3), Inches(0.5),
             desc, font_size=9, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════

out_path = os.path.join(os.path.dirname(__file__), "output", "pLIN_Tool_Architecture.pptx")
os.makedirs(os.path.dirname(out_path), exist_ok=True)
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
